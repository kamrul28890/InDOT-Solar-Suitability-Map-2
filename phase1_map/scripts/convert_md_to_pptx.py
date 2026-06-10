from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

# Simple markdown-to-pptx converter tailored to our slide format
INPUT = Path(r"d:/My Projects/InDOT/MEETING_PRESENTATION.md")
OUTPUT = Path(r"d:/My Projects/InDOT/MEETING_PRESENTATION.pptx")

def parse_slides(md_text):
    parts = [p.strip() for p in md_text.split('\n---\n') if p.strip()]
    slides = []
    for p in parts:
        lines = [l.rstrip() for l in p.splitlines() if l.strip()]
        if not lines:
            continue
        # Title: take first non-empty line, strip leading hashes or headings
        title = lines[0]
        body_lines = lines[1:]
        slides.append((title, body_lines))
    return slides


def add_bullets(shape, bullets):
    tf = shape.text_frame
    tf.clear()
    first = True
    for b in bullets:
        # Use lines starting with '-' as bullets; otherwise preserve as paragraph
        if b.startswith('- '):
            text = b[2:].strip()
            p = tf.add_paragraph()
            p.text = text
            p.level = 0
            p.font.size = Pt(18)
        else:
            # plain paragraph
            p = tf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(14)


def build_pptx(slides):
    prs = Presentation()
    # blank slide layout index may vary; use layout 5 (Title and Content) fallback
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    for title, body in slides:
        slide = prs.slides.add_slide(layout)
        # set title
        try:
            slide.shapes.title.text = title
        except Exception:
            pass
        # body placeholder
        body_shapes = [s for s in slide.shapes if s.is_placeholder]
        if body_shapes:
            shp = body_shapes[1] if len(body_shapes) > 1 else body_shapes[0]
            add_bullets(shp, body)
        else:
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(4.5)
            shp = slide.shapes.add_textbox(left, top, width, height)
            add_bullets(shp, body)
    prs.save(OUTPUT)
    print(f"Saved PPTX to: {OUTPUT}")


def main():
    md = INPUT.read_text(encoding='utf-8')
    slides = parse_slides(md)
    build_pptx(slides)

if __name__ == '__main__':
    main()
